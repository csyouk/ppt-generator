"""
빌더 CLI

사용법:
    python -m src.build \
        --input examples/input.json \
        --config config.yaml \
        --output output.pptx

병합 모드 (여러 JSON 합치기):
    python -m src.build \
        --input chapter1.json chapter2.json chapter3.json \
        --config config.yaml \
        --output full_deck.pptx
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from .design import load_theme
from .layouts import render_slide
from .schema import Deck


def merge_decks(deck_paths: list[Path]) -> Deck:
    """여러 JSON 파일을 하나의 Deck으로 합침"""
    if len(deck_paths) == 1:
        with open(deck_paths[0], encoding="utf-8") as f:
            return Deck.model_validate(json.load(f))

    title = None
    author = None
    slides = []
    for p in deck_paths:
        with open(p, encoding="utf-8") as f:
            data = json.load(f)
        d = Deck.model_validate(data)
        if title is None:
            title = d.title
            author = d.author
        slides.extend(d.slides)
    return Deck(title=title or "Presentation", author=author, slides=slides)


def build(deck: Deck, config_path: Path, output_path: Path, cache_dir: Path):
    theme = load_theme(config_path)

    prs = Presentation()
    prs.slide_width = Inches(theme.layout.slide_w)
    prs.slide_height = Inches(theme.layout.slide_h)

    # 빈 레이아웃 사용 (마스터의 placeholder 영향 제거)
    blank_layout = prs.slide_layouts[6]

    total = len(deck.slides)
    for i, slide_data in enumerate(deck.slides, start=1):
        slide = prs.slides.add_slide(blank_layout)
        try:
            render_slide(slide, slide_data, theme, page_num=i, cache_dir=cache_dir)
        except Exception as e:
            print(
                f"\n[ERROR] 슬라이드 {i} ({slide_data.type}) 렌더링 실패: {e}",
                file=sys.stderr,
            )
            raise
        if i % 10 == 0 or i == total:
            print(f"  [{i:>3}/{total}] {slide_data.type}", file=sys.stderr)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\n✓ 저장됨: {output_path} ({total} slides)", file=sys.stderr)


def main():
    ap = argparse.ArgumentParser(description="JSON → PPTX 빌더")
    ap.add_argument(
        "--input", "-i", nargs="+", required=True,
        help="슬라이드 JSON 파일 (여러 개 가능)",
    )
    ap.add_argument(
        "--config", "-c", default="config.yaml",
        help="브랜드 설정 파일 (기본: config.yaml)",
    )
    ap.add_argument(
        "--output", "-o", default="output.pptx",
        help="출력 파일명",
    )
    ap.add_argument(
        "--cache-dir", default=".cache",
        help="렌더링 캐시 디렉토리",
    )
    args = ap.parse_args()

    deck_paths = [Path(p) for p in args.input]
    for p in deck_paths:
        if not p.exists():
            sys.exit(f"파일 없음: {p}")

    config_path = Path(args.config)
    if not config_path.exists():
        sys.exit(f"설정 파일 없음: {config_path}")

    deck = merge_decks(deck_paths)
    build(
        deck=deck,
        config_path=config_path,
        output_path=Path(args.output),
        cache_dir=Path(args.cache_dir),
    )


if __name__ == "__main__":
    main()
