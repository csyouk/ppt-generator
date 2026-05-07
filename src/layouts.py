"""
슬라이드 레이아웃 — 타입별로 정확히 하나씩.

이 파일이 디자인 일관성의 본체입니다.
- 모든 좌표/크기/색은 theme에서 가져옵니다.
- 같은 타입의 슬라이드는 같은 함수를 거치므로 픽셀 단위로 동일합니다.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from .design import Theme
from .renderers import render_code, render_mermaid, render_table
from .schema import (
    CodeSlide,
    ContentSlide,
    DiagramSlide,
    ImageTextSlide,
    QuoteSlide,
    SectionSlide,
    Slide,
    StatCalloutSlide,
    StepsSlide,
    TableSlide,
    TitleSlide,
    TwoColumnSlide,
)


# ─────────────────────────────────────────────────────────
# 헬퍼 — 텍스트 박스/도형에 디자인 토큰 적용
# ─────────────────────────────────────────────────────────


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_fill(shape, hex_color: str):
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    shape.line.fill.background()


def _add_text(slide, x, y, w, h, text, *, font, size, color,
              bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    # \n으로 분리된 여러 줄 지원
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = _rgb(color)
    return box


def _add_bullets(slide, x, y, w, h, items, *, theme: Theme, size=None):
    """글머리 항목 — 색 있는 점 + 텍스트"""
    if size is None:
        size = theme.sizes.body
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)

        # 컬러 불릿 점
        bullet_run = p.add_run()
        bullet_run.text = "● "
        bullet_run.font.name = theme.fonts.body
        bullet_run.font.size = Pt(size)
        bullet_run.font.color.rgb = _rgb(theme.colors.accent)

        # 본문
        text_run = p.add_run()
        text_run.text = item
        text_run.font.name = theme.fonts.body
        text_run.font.size = Pt(size)
        text_run.font.color.rgb = _rgb(theme.colors.ink)
    return box


def _add_image_fit(slide, image_path: Path, x, y, w, h):
    """
    이미지를 (x,y)-(x+w, y+h) 영역에 비율 유지하며 배치.
    """
    from PIL import Image as PILImage

    im = PILImage.open(image_path)
    iw, ih = im.size
    aspect = iw / ih
    box_aspect = w / h

    if aspect > box_aspect:
        # 너비 기준
        draw_w = w
        draw_h = w / aspect
        draw_x = x
        draw_y = y + (h - draw_h) / 2
    else:
        # 높이 기준
        draw_h = h
        draw_w = h * aspect
        draw_x = x + (w - draw_w) / 2
        draw_y = y

    slide.shapes.add_picture(
        str(image_path),
        Inches(draw_x), Inches(draw_y),
        Inches(draw_w), Inches(draw_h),
    )


# ─────────────────────────────────────────────────────────
# 공통 요소 — 배경, footer, 페이지 번호
# ─────────────────────────────────────────────────────────


def _paint_background(slide, theme: Theme, color_hex: str):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(theme.layout.slide_w), Inches(theme.layout.slide_h),
    )
    _set_fill(bg, color_hex)
    bg.shadow.inherit = False
    # 배경은 가장 뒤로
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def _add_footer(slide, theme: Theme, page_num: int | None = None):
    if not theme.show_footer:
        return
    L = theme.layout
    if theme.brand_name:
        _add_text(
            slide, L.margin_x, L.footer_y,
            L.slide_w - 2 * L.margin_x, L.footer_h,
            theme.brand_name,
            font=theme.fonts.body, size=theme.sizes.caption,
            color=theme.colors.muted,
            align=PP_ALIGN.LEFT,
        )
    if theme.show_page_numbers and page_num is not None:
        _add_text(
            slide, L.margin_x, L.footer_y,
            L.slide_w - 2 * L.margin_x, L.footer_h,
            str(page_num),
            font=theme.fonts.body, size=theme.sizes.caption,
            color=theme.colors.muted,
            align=PP_ALIGN.RIGHT,
        )


def _add_title_block(slide, theme: Theme, title: str, subtitle: str | None = None,
                     eyebrow: str | None = None):
    """슬라이드 상단 타이틀 영역 — title, subtitle, eyebrow 배치"""
    L = theme.layout
    S = theme.sizes
    y = L.title_y

    if eyebrow:
        _add_text(
            slide, L.margin_x, y, L.slide_w - 2 * L.margin_x, 0.3,
            eyebrow.upper(),
            font=theme.fonts.heading, size=S.eyebrow,
            color=theme.colors.accent, bold=True,
        )
        y += 0.35

    _add_text(
        slide, L.margin_x, y, L.slide_w - 2 * L.margin_x, L.title_h,
        title,
        font=theme.fonts.heading, size=S.title_lg,
        color=theme.colors.primary, bold=True,
    )

    if subtitle:
        # 타이틀 박스 바로 아래에 subtitle (겹치지 않게)
        _add_text(
            slide, L.margin_x, y + L.title_h + 0.05,
            L.slide_w - 2 * L.margin_x, 0.4,
            subtitle,
            font=theme.fonts.body, size=S.body,
            color=theme.colors.muted,
        )


# ─────────────────────────────────────────────────────────
# 레이아웃 — 슬라이드 타입별
# ─────────────────────────────────────────────────────────


def layout_title(slide, data: TitleSlide, theme: Theme, page_num: int):
    """표지 — 다크 배경, 큰 타이틀 중앙"""
    _paint_background(slide, theme, theme.colors.primary)
    L, S = theme.layout, theme.sizes

    cy = L.slide_h / 2 - 1.0  # 중앙보다 약간 위

    if data.eyebrow:
        _add_text(
            slide, L.margin_x, cy, L.slide_w - 2 * L.margin_x, 0.4,
            data.eyebrow.upper(),
            font=theme.fonts.heading, size=S.eyebrow + 2,
            color=theme.colors.accent, bold=True, align=PP_ALIGN.CENTER,
        )
        cy += 0.5

    _add_text(
        slide, L.margin_x, cy, L.slide_w - 2 * L.margin_x, 1.6,
        data.title,
        font=theme.fonts.heading, size=S.title_xl,
        color=theme.colors.on_primary, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    if data.subtitle:
        _add_text(
            slide, L.margin_x, cy + 1.6,
            L.slide_w - 2 * L.margin_x, 0.6,
            data.subtitle,
            font=theme.fonts.body, size=S.body_lg,
            color=theme.colors.on_primary,
            align=PP_ALIGN.CENTER,
        )


def layout_section(slide, data: SectionSlide, theme: Theme, page_num: int):
    """섹션 구분 — 다크 배경 + 큰 번호"""
    _paint_background(slide, theme, theme.colors.primary)
    L, S = theme.layout, theme.sizes

    # 큰 번호 (위쪽)
    _add_text(
        slide, L.margin_x, 1.4, L.slide_w - 2 * L.margin_x, 2.0,
        data.section_number,
        font=theme.fonts.heading, size=110,
        color=theme.colors.accent, bold=True,
    )

    # 타이틀 (큰 번호와 충분한 간격)
    _add_text(
        slide, L.margin_x, 4.0, L.slide_w - 2 * L.margin_x, 1.0,
        data.title,
        font=theme.fonts.heading, size=S.title_lg,
        color=theme.colors.on_primary, bold=True,
    )

    # 설명
    if data.description:
        _add_text(
            slide, L.margin_x, 5.1, L.slide_w - 2 * L.margin_x, 1.5,
            data.description,
            font=theme.fonts.body, size=S.body,
            color=theme.colors.on_primary,
        )


def layout_content(slide, data: ContentSlide, theme: Theme, page_num: int):
    """일반 컨텐츠"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L, S = theme.layout, theme.sizes
    if data.bullets:
        _add_bullets(
            slide, L.margin_x, L.content_y,
            L.slide_w - 2 * L.margin_x, L.content_h,
            data.bullets, theme=theme, size=S.body_lg,
        )
    elif data.paragraph:
        _add_text(
            slide, L.margin_x, L.content_y,
            L.slide_w - 2 * L.margin_x, L.content_h,
            data.paragraph,
            font=theme.fonts.body, size=S.body_lg,
            color=theme.colors.ink,
        )

    if data.footer_note:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.footer_note,
            font=theme.fonts.body, size=S.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_two_column(slide, data: TwoColumnSlide, theme: Theme, page_num: int):
    """2단 비교"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title)
    _add_footer(slide, theme, page_num)

    L, S = theme.layout, theme.sizes
    col_gap = 0.4
    col_w = (L.slide_w - 2 * L.margin_x - col_gap) / 2

    for i, (heading, bullets, x) in enumerate([
        (data.left_heading, data.left_bullets, L.margin_x),
        (data.right_heading, data.right_bullets, L.margin_x + col_w + col_gap),
    ]):
        # 컬럼 헤딩 (컬러 박스)
        head_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(L.content_y),
            Inches(col_w), Inches(0.5),
        )
        _set_fill(head_box, theme.colors.secondary)
        head_tf = head_box.text_frame
        head_tf.margin_left = Inches(0.2)
        head_tf.margin_right = Inches(0.2)
        head_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        head_p = head_tf.paragraphs[0]
        head_run = head_p.add_run()
        head_run.text = heading
        head_run.font.name = theme.fonts.heading
        head_run.font.size = Pt(S.title_md - 4)
        head_run.font.bold = True
        head_run.font.color.rgb = _rgb(theme.colors.primary)

        # 글머리
        _add_bullets(
            slide, x + 0.05, L.content_y + 0.7,
            col_w - 0.1, L.content_h - 0.7,
            bullets, theme=theme, size=S.body,
        )


def layout_table(slide, data: TableSlide, theme: Theme, page_num: int,
                 cache_dir: Path):
    """Markdown 테이블 — 이미지로 렌더링"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L = theme.layout
    img = render_table(data.markdown_table, theme, cache_dir)
    _add_image_fit(
        slide, img,
        L.margin_x, L.content_y,
        L.slide_w - 2 * L.margin_x,
        L.content_h - (0.4 if data.caption else 0),
    )

    if data.caption:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.caption,
            font=theme.fonts.body, size=theme.sizes.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_diagram(slide, data: DiagramSlide, theme: Theme, page_num: int,
                   cache_dir: Path):
    """Mermaid 다이어그램"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L = theme.layout
    img = render_mermaid(data.mermaid_code, theme, cache_dir)
    _add_image_fit(
        slide, img,
        L.margin_x, L.content_y,
        L.slide_w - 2 * L.margin_x,
        L.content_h - (0.4 if data.caption else 0),
    )

    if data.caption:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.caption,
            font=theme.fonts.body, size=theme.sizes.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_code(slide, data: CodeSlide, theme: Theme, page_num: int,
                cache_dir: Path):
    """코드 슬라이드 — syntax highlight"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L = theme.layout
    img = render_code(data.code, data.language, theme, cache_dir)
    _add_image_fit(
        slide, img,
        L.margin_x, L.content_y,
        L.slide_w - 2 * L.margin_x,
        L.content_h - (0.4 if data.caption else 0),
    )

    if data.caption:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.caption,
            font=theme.fonts.body, size=theme.sizes.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_quote(slide, data: QuoteSlide, theme: Theme, page_num: int):
    """인용 슬라이드"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_footer(slide, theme, page_num)
    L, S = theme.layout, theme.sizes

    # 큰 따옴표 장식 — 위쪽에 단독 배치
    _add_text(
        slide, L.margin_x + 0.4, 1.6, 1.2, 1.4,
        "\u201C",
        font=theme.fonts.heading, size=140,
        color=theme.colors.accent, bold=True,
    )

    # 인용 본문 — 따옴표 아래, 좌측 정렬
    _add_text(
        slide, L.margin_x + 0.4, 3.2,
        L.slide_w - 2 * L.margin_x - 0.8, 2.6,
        data.quote,
        font=theme.fonts.heading, size=S.title_md,
        color=theme.colors.ink, italic=True,
    )

    # 출처
    if data.attribution:
        _add_text(
            slide, L.margin_x + 0.4, 5.9,
            L.slide_w - 2 * L.margin_x - 0.8, 0.4,
            f"— {data.attribution}",
            font=theme.fonts.body, size=S.body,
            color=theme.colors.muted,
        )


def layout_stat_callout(slide, data: StatCalloutSlide, theme: Theme, page_num: int):
    """큰 숫자 강조 — 1~3개 stat을 나란히"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L, S = theme.layout, theme.sizes
    n = len(data.stats)
    if n == 0:
        return
    n = min(n, 3)  # 최대 3개

    # 사용 가능한 영역
    area_w = L.slide_w - 2 * L.margin_x
    gap = 0.4
    col_w = (area_w - gap * (n - 1)) / n
    cy = L.content_y + 0.3

    for i, stat in enumerate(data.stats[:n]):
        x = L.margin_x + i * (col_w + gap)

        # 큰 숫자
        _add_text(
            slide, x, cy, col_w, 1.6,
            stat.value,
            font=theme.fonts.heading, size=72,
            color=theme.colors.accent, bold=True,
            align=PP_ALIGN.CENTER,
        )

        # 라벨
        _add_text(
            slide, x, cy + 1.7, col_w, 0.5,
            stat.label,
            font=theme.fonts.heading, size=S.body_lg,
            color=theme.colors.primary, bold=True,
            align=PP_ALIGN.CENTER,
        )

        # 부연 (선택)
        if stat.description:
            _add_text(
                slide, x, cy + 2.3, col_w, 1.2,
                stat.description,
                font=theme.fonts.body, size=S.body_sm,
                color=theme.colors.muted,
                align=PP_ALIGN.CENTER,
            )

    # footer note
    if data.footer_note:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.footer_note,
            font=theme.fonts.body, size=S.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_image_text(slide, data: ImageTextSlide, theme: Theme, page_num: int):
    """이미지 + 텍스트 좌우 분할"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L, S = theme.layout, theme.sizes
    gap = 0.4
    area_w = L.slide_w - 2 * L.margin_x
    half_w = (area_w - gap) / 2

    if data.image_side == "left":
        img_x = L.margin_x
        text_x = L.margin_x + half_w + gap
    else:
        text_x = L.margin_x
        img_x = L.margin_x + half_w + gap

    # 이미지 영역
    image_path = Path(data.image_path)
    if not image_path.is_absolute():
        # 상대경로면 cwd 기준
        image_path = Path.cwd() / image_path

    if image_path.exists():
        _add_image_fit(
            slide, image_path,
            img_x, L.content_y,
            half_w, L.content_h - 0.5,
        )
    else:
        # 이미지 없으면 placeholder
        ph = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(img_x), Inches(L.content_y),
            Inches(half_w), Inches(L.content_h - 0.5),
        )
        _set_fill(ph, theme.colors.secondary)
        ph_tf = ph.text_frame
        ph_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        ph_p = ph_tf.paragraphs[0]
        ph_p.alignment = PP_ALIGN.CENTER
        ph_run = ph_p.add_run()
        ph_run.text = f"[이미지 없음]\n{data.image_path}"
        ph_run.font.name = theme.fonts.body
        ph_run.font.size = Pt(S.caption)
        ph_run.font.color.rgb = _rgb(theme.colors.muted)

    # 텍스트 영역
    if data.bullets:
        _add_bullets(
            slide, text_x, L.content_y,
            half_w, L.content_h - 0.5,
            data.bullets, theme=theme, size=S.body,
        )
    elif data.paragraph:
        _add_text(
            slide, text_x, L.content_y, half_w, L.content_h - 0.5,
            data.paragraph,
            font=theme.fonts.body, size=S.body,
            color=theme.colors.ink,
        )

    # 캡션
    if data.caption:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.caption,
            font=theme.fonts.body, size=S.caption,
            color=theme.colors.muted, italic=True,
        )


def layout_steps(slide, data: StepsSlide, theme: Theme, page_num: int):
    """단계별 흐름 — 번호 매긴 카드 3~5개"""
    _paint_background(slide, theme, theme.colors.paper)
    _add_title_block(slide, theme, data.title, data.subtitle)
    _add_footer(slide, theme, page_num)

    L, S = theme.layout, theme.sizes
    n = len(data.steps)
    if n == 0:
        return
    n = min(n, 5)

    area_w = L.slide_w - 2 * L.margin_x
    gap = 0.25
    col_w = (area_w - gap * (n - 1)) / n
    cy = L.content_y + 0.2
    card_h = 3.2

    for i, step in enumerate(data.steps[:n]):
        x = L.margin_x + i * (col_w + gap)

        # 번호 원
        circle_size = 0.8
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + (col_w - circle_size) / 2),
            Inches(cy),
            Inches(circle_size), Inches(circle_size),
        )
        _set_fill(circle, theme.colors.accent)
        c_tf = circle.text_frame
        c_tf.margin_left = Inches(0)
        c_tf.margin_right = Inches(0)
        c_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        c_p = c_tf.paragraphs[0]
        c_p.alignment = PP_ALIGN.CENTER
        c_run = c_p.add_run()
        c_run.text = str(i + 1)
        c_run.font.name = theme.fonts.heading
        c_run.font.size = Pt(28)
        c_run.font.bold = True
        c_run.font.color.rgb = _rgb("FFFFFF")

        # 단계 제목
        _add_text(
            slide, x, cy + circle_size + 0.2, col_w, 0.6,
            step.title,
            font=theme.fonts.heading, size=S.title_md - 6,
            color=theme.colors.primary, bold=True,
            align=PP_ALIGN.CENTER,
        )

        # 단계 설명
        _add_text(
            slide, x + 0.1, cy + circle_size + 0.85,
            col_w - 0.2, card_h - circle_size - 0.85,
            step.description,
            font=theme.fonts.body, size=S.body_sm,
            color=theme.colors.ink,
            align=PP_ALIGN.CENTER,
        )

    # footer
    if data.footer_note:
        _add_text(
            slide, L.margin_x, L.footer_y - 0.4,
            L.slide_w - 2 * L.margin_x, 0.3,
            data.footer_note,
            font=theme.fonts.body, size=S.caption,
            color=theme.colors.muted, italic=True,
        )


# ─────────────────────────────────────────────────────────
# 디스패처
# ─────────────────────────────────────────────────────────


def render_slide(slide, data: Slide, theme: Theme, page_num: int,
                 cache_dir: Path):
    if isinstance(data, TitleSlide):
        layout_title(slide, data, theme, page_num)
    elif isinstance(data, SectionSlide):
        layout_section(slide, data, theme, page_num)
    elif isinstance(data, ContentSlide):
        layout_content(slide, data, theme, page_num)
    elif isinstance(data, TwoColumnSlide):
        layout_two_column(slide, data, theme, page_num)
    elif isinstance(data, TableSlide):
        layout_table(slide, data, theme, page_num, cache_dir)
    elif isinstance(data, DiagramSlide):
        layout_diagram(slide, data, theme, page_num, cache_dir)
    elif isinstance(data, CodeSlide):
        layout_code(slide, data, theme, page_num, cache_dir)
    elif isinstance(data, QuoteSlide):
        layout_quote(slide, data, theme, page_num)
    elif isinstance(data, StatCalloutSlide):
        layout_stat_callout(slide, data, theme, page_num)
    elif isinstance(data, ImageTextSlide):
        layout_image_text(slide, data, theme, page_num)
    elif isinstance(data, StepsSlide):
        layout_steps(slide, data, theme, page_num)
    else:
        raise ValueError(f"Unknown slide type: {type(data)}")
